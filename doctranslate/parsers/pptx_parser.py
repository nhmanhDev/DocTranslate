import os
from pptx import Presentation
from doctranslate.schema import UnifiedDocument, DocumentBlock, Span

def parse_pptx(file_path: str) -> UnifiedDocument:
    """Parse a PPTX file into a UnifiedDocument model."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPTX file not found: {file_path}")

    prs = Presentation(file_path)
    unified_doc = UnifiedDocument(metadata={"title": os.path.basename(file_path), "format": "pptx"})

    block_counter = 0

    def parse_text_frame_to_blocks(tf, shape_id_prefix):
        nonlocal block_counter
        blocks = []
        for p_idx, p in enumerate(tf.paragraphs):
            content = p.text.strip()
            if not content:
                continue

            block_counter += 1
            spans = []
            for run in p.runs:
                color = None
                if run.font.color and run.font.color.rgb:
                    color = str(run.font.color.rgb)

                span = Span(
                    text=run.text,
                    bold=bool(run.font.bold),
                    italic=bool(run.font.italic),
                    underline=bool(run.font.underline),
                    font_size=run.font.size.pt if run.font.size else None,
                    font_name=run.font.name,
                    color=color,
                )
                spans.append(span)

            block = DocumentBlock(
                id=f"{shape_id_prefix}_p_{block_counter}",
                type="paragraph",
                content=content,
                spans=spans,
                attributes={"slide_index": slide_idx + 1}
            )
            blocks.append(block)
        return blocks

    # Iterate slides
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            shape_prefix = f"s_{slide_idx}_sh_{shape_idx}"
            
            # 1. Text frame shapes
            if shape.has_text_frame:
                blocks = parse_text_frame_to_blocks(shape.text_frame, shape_prefix)
                unified_doc.blocks.extend(blocks)
                
            # 2. Table shapes
            elif shape.has_table:
                block_counter += 1
                table_rows = []
                table_obj = shape.table
                for row_idx, row in enumerate(table_obj.rows):
                    row_cells = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_prefix = f"{shape_prefix}_t_{row_idx}_{col_idx}"
                        cell_blocks = parse_text_frame_to_blocks(cell.text_frame, cell_prefix)
                        row_cells.append(cell_blocks)
                    table_rows.append(row_cells)

                # Add table block
                unified_doc.blocks.append(DocumentBlock(
                    id=f"{shape_prefix}_table_{block_counter}",
                    type="table",
                    content="",
                    table_rows=table_rows,
                    attributes={"slide_index": slide_idx + 1, "rows": len(table_obj.rows), "cols": len(table_obj.columns)}
                ))

    return unified_doc

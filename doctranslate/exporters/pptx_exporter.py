import shutil
import os
from pptx import Presentation
from doctranslate.schema import UnifiedDocument

def export_pptx(
    template_path: str,
    output_path: str,
    translated_doc: UnifiedDocument,
    mode: str = "mono"  # "mono" or "dual"
):
    """
    Export the translated document to PPTX.
    Uses the original file as a template to preserve layouts and slide themes.
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template PPTX file not found: {template_path}")

    # Copy template to output path
    shutil.copy2(template_path, output_path)
    prs = Presentation(output_path)

    # Collect blocks lookup
    blocks_by_id = {}
    def collect_blocks(blocks):
        for block in blocks:
            blocks_by_id[block.id] = block
            if block.table_rows:
                for row in block.table_rows:
                    for cell in row:
                        collect_blocks(cell)
                        
    collect_blocks(translated_doc.blocks)

    block_counter = 0

    def apply_translation_to_text_frame(tf, shape_id_prefix):
        nonlocal block_counter
        for p_idx, p in enumerate(tf.paragraphs):
            # We must verify if this paragraph had content and is in our block list
            content = p.text.strip()
            if not content:
                continue

            block_counter += 1
            block_id = f"{shape_id_prefix}_p_{block_counter}"
            block = blocks_by_id.get(block_id)
            if not block or not block.translated_content:
                continue

            translated_text = block.translated_content

            if mode == "mono":
                # Clear all runs and set the translation on the first run
                if p.runs:
                    p.runs[0].text = translated_text
                    for run in p.runs[1:]:
                        run.text = ""
                else:
                    p.text = translated_text
            elif mode == "dual":
                # In presentation slides, space is limited, but we can append dual text with a separator
                # We add a newline run and the translated text run
                p.add_run()
                p.runs[-1].text = "\n"
                trans_run = p.add_run()
                trans_run.text = translated_text
                trans_run.font.italic = True
                try:
                    from pptx.dml.color import RGBColor
                    trans_run.font.color.rgb = RGBColor(128, 128, 128)
                except Exception:
                    pass

    # Traverse slides and shapes
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            shape_prefix = f"s_{slide_idx}_sh_{shape_idx}"
            
            if shape.has_text_frame:
                apply_translation_to_text_frame(shape.text_frame, shape_prefix)
            elif shape.has_table:
                block_counter += 1
                table_obj = shape.table
                for row_idx, row in enumerate(table_obj.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell_prefix = f"{shape_prefix}_t_{row_idx}_{col_idx}"
                        apply_translation_to_text_frame(cell.text_frame, cell_prefix)

    prs.save(output_path)
